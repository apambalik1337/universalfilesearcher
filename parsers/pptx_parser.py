"""
pptx_parser.py

Parses PowerPoint files (.pptx) using python-pptx.

Searches slide titles, text boxes/paragraphs, tables, and speaker notes.
Results preserve the slide number and label which part of the slide
matched (Title / Text box / Table / Speaker notes).
"""

from pptx import Presentation

from .base_parser import BaseParser, ParseResult, Record


class PptxParser(BaseParser):
    extensions = (".pptx",)
    display_name = "PowerPoint (.pptx)"

    def parse(self, file_path: str) -> ParseResult:
        try:
            prs = Presentation(file_path)
        except Exception as e:
            return ParseResult(success=False, error=f"Failed to open presentation: {e}")

        records = []
        try:
            for slide_idx, slide in enumerate(prs.slides, start=1):
                title_text = None
                title_shape_id = None
                try:
                    if slide.shapes.title is not None:
                        title_text = (slide.shapes.title.text or "").strip()
                        title_shape_id = slide.shapes.title.shape_id
                except Exception:
                    title_text = None

                if title_text:
                    records.append(Record(content=title_text, slide=slide_idx, location="Title"))

                for shape in slide.shapes:
                    is_title_shape = title_shape_id is not None and getattr(shape, "shape_id", None) == title_shape_id

                    if getattr(shape, "has_text_frame", False) and not is_title_shape:
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in para.runs).strip() or para.text.strip()
                            if text:
                                records.append(Record(content=text, slide=slide_idx, location="Text box"))

                    if getattr(shape, "has_table", False):
                        table = shape.table
                        for r_idx, row in enumerate(table.rows, start=1):
                            for c_idx, cell in enumerate(row.cells, start=1):
                                text = (cell.text or "").strip()
                                if text:
                                    records.append(
                                        Record(
                                            content=text,
                                            slide=slide_idx,
                                            location=f"Table, Row {r_idx}, Column {c_idx}",
                                        )
                                    )

                try:
                    if slide.has_notes_slide:
                        notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                        if notes_text:
                            records.append(Record(content=notes_text, slide=slide_idx, location="Speaker notes"))
                except Exception:
                    pass
        except Exception as e:
            return ParseResult(success=False, error=f"Error while reading slide contents: {e}", records=records)

        if not records:
            return ParseResult(success=True, records=[], warning="No content found in presentation")

        return ParseResult(success=True, records=records)

    def is_dependency_available(self) -> bool:
        try:
            import pptx  # noqa
            return True
        except ImportError:
            return False

    def dependency_message(self) -> str:
        return "Install with: pip install python-pptx"
